#!/usr/bin/env python3
import os, sys, json, time, hashlib, logging
from pathlib import Path
from typing import List, Dict, Tuple, Optional
os.environ["user_agent"] = "cohere-langchain-client"

# --- Third-party ---
import numpy as np
import oracledb
import oci
from dotenv import load_dotenv
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import chardet
#from langchain_community.embeddings import OCIGenAIEmbeddings
from langchain_cohere import CohereEmbeddings



# ----------------- LOGGING -----------------
logging.basicConfig(
    level=os.getenv("LOG_LEVEL", "INFO"),
    format="%(asctime)s %(levelname)s %(message)s",
)
log = logging.getLogger("ingest")

# ----------------- ENV/CONFIG -----------------
load_dotenv()  # reads .env if present

# Required for embeddings
COMPARTMENT_OCID = os.environ["COMPARTMENT_OCID"]
OCI_REGION       = os.environ["OCI_REGION"]

# ADB connection (wallet thick-mode)
TNS_ADMIN   = os.environ["TNS_ADMIN"]           # e.g. /home/<you>/wallet
ADB_USER    = os.environ["ADB_USER"]            # e.g. ADMIN
ADB_PASSWORD= os.environ["ADB_PASSWORD"]
ADB_DSN     = os.environ["ADB_DSN"]             # e.g. aitutordb_tp
WALLET_PW   = os.getenv("WALLET_PASSWORD")      # only if you set one at wallet download

# Ingest config
COURSE      = os.getenv("COURSE", "ALISHA_DEV")
BUCKET      = os.getenv("BUCKET", "course-docs")
SKIP_OCI    = "0" #os.getenv("SKIP_OCI", "0") == "0"
RAW_DIR     = Path(os.getenv("RAW_DIR", "course_docs"))
RAW_DIR.mkdir(parents=True, exist_ok=True)

EMBED_MODEL_ID = os.getenv("EMBED_MODEL_ID", "cohere.embed-english-v3.0")
EMBED_DIM      = int(os.getenv("EMBED_DIM", "1024"))

BATCH_SIZE      = int(os.getenv("BATCH_SIZE", "64"))
COMMIT_EVERY_N  = int(os.getenv("COMMIT_EVERY_N", "512"))

print("✅ Using python-oracledb in THIN mode with wallet.")
print(f"Wallet path set via TNS_ADMIN={TNS_ADMIN}")

# ----------------- OCI OBJECT STORAGE -----------------
def make_os_client():
    cfg_path = Path("~/.oci/config").expanduser()
    if cfg_path.exists():
        config = oci.config.from_file(str(cfg_path), "DEFAULT")
        return oci.object_storage.ObjectStorageClient(config)
    signer = oci.auth.signers.get_resource_principals_signer()
    return oci.object_storage.ObjectStorageClient(config={}, signer=signer)

def download_bucket_to(dir_path: Path, bucket: str) -> None:
    if SKIP_OCI:
        log.info("[objstore] SKIP_OCI=1 → skipping object storage download")
        return
    os_client = make_os_client()
    namespace = os_client.get_namespace().data
    objs = os_client.list_objects(namespace, bucket).data.objects
    for obj in objs:
        dest = dir_path / obj.name
        dest.parent.mkdir(parents=True, exist_ok=True)
        log.info("↓ %s", obj.name)
        with open(dest, "wb") as f:
            data = os_client.get_object(namespace, bucket, obj.name).data
            f.write(data.content)

# ----------------- SIMPLE PARSERS -----------------
def read_txt(path: Path) -> Tuple[str, dict]:
    raw = path.read_bytes()
    enc = chardet.detect(raw).get("encoding") or "utf-8"
    return raw.decode(enc, errors="ignore"), {}

def read_pdf(path: Path) -> Tuple[str, dict]:
    reader = PdfReader(str(path))
    pages = []
    for i, pg in enumerate(reader.pages):
        try:
            t = pg.extract_text() or ""
        except Exception:
            t = ""
        pages.append(t)
    return "\n\n".join(pages), {"page_count": len(pages)}

def read_docx(path: Path) -> Tuple[str, dict]:
    doc = DocxDocument(str(path))
    paras = [p.text for p in doc.paragraphs if p.text]
    return "\n\n".join(paras), {}

def read_pptx(path: Path) -> Tuple[str, dict]:
    prs = Presentation(str(path))
    slides_txt: List[str] = []
    for s in prs.slides:
        parts = []
        for shp in s.shapes:
            if hasattr(shp, "text"):
                parts.append(shp.text)
        slides_txt.append("\n".join(parts))
    return "\n\n".join(slides_txt), {"slide_count": len(prs.slides)}

def load_text(path: Path) -> Tuple[str, dict]:
    ext = path.suffix.lower()
    if ext == ".pdf":  return read_pdf(path)
    if ext == ".docx": return read_docx(path)
    if ext == ".pptx": return read_pptx(path)
    if ext in {".txt", ".md"}: return read_txt(path)
    raise ValueError(f"Unsupported file type: {ext}")

# ----------------- CHUNKING -----------------
def window_chunks(text: str, max_chars=1500, overlap=300) -> List[str]:
    if not text: return []
    out, start, n = [], 0, len(text)
    while start < n:
        end = min(start + max_chars, n)
        out.append(text[start:end])
        if end == n: break
        start = max(0, end - overlap)
    return out

def header_aware_chunks(text: str, max_chars=1500, overlap=300) -> List[str]:
    import re
    if not text: return []
    paras = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]

    def is_header(p: str) -> bool:
        line = p.splitlines()[0] if p else ""
        if len(line) <= 80 and (line.isupper() or line.endswith(":")): return True
        return bool(re.match(r"^(\d+(\.\d+)*|[IVXLC]+\.)\s+\S", line))

    chunks, buf = [], ""
    for p in paras:
        add = ("\n\n" + p) if buf else p
        if len(buf + add) <= max_chars:
            buf = buf + add
        else:
            if buf: chunks.append(buf)
            tail = buf[-overlap:] if overlap and buf else ""
            buf = (tail + "\n\n" + p) if tail else p
        if is_header(p) and len(buf) > max_chars:
            chunks.append(buf); buf = ""
    if buf: chunks.append(buf)

    final = []
    for c in chunks:
        final.extend(window_chunks(c, max_chars, overlap) if len(c) > max_chars else [c])
    return final

def make_chunks(text: str, strategy="hybrid", max_chars=1500, overlap=300) -> List[str]:
    return header_aware_chunks(text, max_chars, overlap) if strategy != "window" else window_chunks(text, max_chars, overlap)

# ----------------- EMBEDDINGS (OCI GenAI) -----------------
"""def make_oci_embedder():
    endpoint = f"https://inference.generativeai.{OCI_REGION}.oci.oraclecloud.com"
    return OCIGenAIEmbeddings(
        compartment_id=COMPARTMENT_OCID,
        model_id=EMBED_MODEL_ID,
        service_endpoint=endpoint
    )"""

def make_oci_embedder():
    from langchain_cohere import CohereEmbeddings
    api_key = os.getenv("COHERE_API_KEY")
    return CohereEmbeddings(model="embed-english-v3.0", cohere_api_key=api_key)

def embed_texts(embedder, texts: List[str]) -> List[List[float]]:
    # LangChain wrapper handles batching internally; you can still split if your docs are huge.
    return embedder.embed_documents(texts)

def connect_adb():
    user = os.getenv("ADB_USER")
    pwd  = os.getenv("ADB_PASSWORD")
    dsn  = os.getenv("ADB_DSN")
    cfg  = os.getenv("TNS_ADMIN")

    print(f"🔐 Using wallet from: {cfg}")
    print(f"➡ Connecting using FULL DSN: {dsn}")

    return oracledb.connect(
        user=user,
        password=pwd,
        dsn=dsn,
        config_dir=cfg,          # ✅ Explicit wallet directory
        wallet_location=cfg,     # ✅ Ensures mutual TLS uses wallet
        ssl_server_dn_match=True
    )


# ----------------- DB PREP (IDEMPOTENT) -----------------
def ensure_unique_ingest_index(conn):
    """Unique constraint to prevent duplicate inserts: (course, source_uri, chunk_no)."""
    sql = """
    BEGIN
      EXECUTE IMMEDIATE 'CREATE UNIQUE INDEX doc_chunks_unique_ingest
                         ON doc_chunks(course, source_uri, chunk_no)';
    EXCEPTION
      WHEN OTHERS THEN
        IF SQLCODE != -955 THEN RAISE; END IF; -- already exists
    END;
    """
    with conn.cursor() as cur:
        cur.execute(sql)
    conn.commit()

# ----------------- INSERT (UPSERT) -----------------
def upsert_chunks(conn, rows: List[Dict]):
    """
    Idempotent MERGE on (course, source_uri, chunk_no).
    embedding bound as VECTOR; metadata as JSON.
    """
    # normalize types
    for r in rows:
        r["embedding"] = [float(x) for x in r["embedding"]]
        r["metadata"]  = json.dumps(r.get("metadata", {}))

    merge_sql = """
    MERGE INTO doc_chunks d
    USING (SELECT :course      AS course,
                  :source_uri  AS source_uri,
                  :chunk_no    AS chunk_no,
                  :section     AS section,
                  :page_from   AS page_from,
                  :page_to     AS page_to,
                  :content     AS content,
                  :embedding   AS embedding,
                  :metadata    AS metadata
           FROM dual) s
    ON (d.course = s.course AND d.source_uri = s.source_uri AND d.chunk_no = s.chunk_no)
    WHEN MATCHED THEN UPDATE SET
         d.section = s.section,
         d.page_from = s.page_from,
         d.page_to   = s.page_to,
         d.content   = s.content,
         d.embedding = s.embedding,
         d.metadata  = s.metadata
    WHEN NOT MATCHED THEN INSERT
         (course, source_uri, chunk_no, section, page_from, page_to, content, embedding, metadata)
         VALUES
         (s.course, s.source_uri, s.chunk_no, s.section, s.page_from, s.page_to, s.content, s.embedding, s.metadata)
    """

    with conn.cursor() as cur:
        cur.setinputsizes(
            course=oracledb.DB_TYPE_VARCHAR,
            source_uri=oracledb.DB_TYPE_VARCHAR,
            chunk_no=oracledb.DB_TYPE_NUMBER,
            section=oracledb.DB_TYPE_VARCHAR,
            page_from=oracledb.DB_TYPE_NUMBER,
            page_to=oracledb.DB_TYPE_NUMBER,
            content=oracledb.DB_TYPE_CLOB,
            embedding=oracledb.DB_TYPE_VECTOR,
            metadata=oracledb.DB_TYPE_JSON
        )

        # execute in manageable batches, commit periodically
        for i in range(0, len(rows), BATCH_SIZE):
            batch = rows[i:i+BATCH_SIZE]
            cur.executemany(merge_sql, batch)
            if (i // BATCH_SIZE) % max(1, (COMMIT_EVERY_N // BATCH_SIZE)) == 0:
                conn.commit()
        conn.commit()

# ----------------- MAIN -----------------
def ingest_paths(paths: List[Path]):
    if not paths:
        log.warning("No files to ingest")
        return

    embedder = make_oci_embedder()
    to_upsert: List[Dict] = []

    for p in paths:
        if not p.is_file():
            continue
        if p.suffix.lower() not in {".pdf", ".docx", ".pptx", ".txt", ".md"}:
            continue

        log.info("[parse] %s", p.name)
        text, meta = load_text(p)
        chunks = make_chunks(text, strategy="hybrid", max_chars=1500, overlap=300)
        page_to = meta.get("page_count") or meta.get("slide_count")
        for i, c in enumerate(chunks, start=1):
            to_upsert.append({
                "course": COURSE,
                "source_uri": f"oci://{BUCKET}/{p.name}" if not SKIP_OCI else f"local://{p.name}",
                "chunk_no": i,
                "section": None,
                "page_from": 1 if page_to else None,
                "page_to": page_to,
                "content": c,
                "metadata": {"strategy":"hybrid","bucket":BUCKET,"file":p.name}
            })

    log.info("[chunks] total prepared: %d", len(to_upsert))
    if not to_upsert:
        return

    # embed in batches
    for i in range(0, len(to_upsert), BATCH_SIZE):
        batch = to_upsert[i:i+BATCH_SIZE]
        texts = [r["content"] for r in batch]
        vecs  = embed_texts(embedder, texts)
        for r, v in zip(batch, vecs):
            r["embedding"] = v

    # insert/upsert
    conn = connect_adb()
    try:
        try:
            ensure_unique_ingest_index(conn)
        except Exception as e:
            log.warning(f"[index] Skipping index creation (may already exist): {e}")

        upsert_chunks(conn, to_upsert)
    finally:
        conn.close()

    log.info("[done] upserted %d rows into DOC_CHUNKS", len(to_upsert))

def main():
    # 1) optionally sync from Object Storage
    download_bucket_to(RAW_DIR, BUCKET)

    # 2) if a specific file is given, ingest just that; otherwise ingest all under RAW_DIR
    if len(sys.argv) > 1:
        paths = [Path(sys.argv[1])]
    else:
        paths = [p for p in sorted(RAW_DIR.rglob("*")) if p.is_file()]

    ingest_paths(paths)

if __name__ == "__main__":
    main()