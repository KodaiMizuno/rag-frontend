import os, json, hashlib
from pathlib import Path
from typing import Tuple, List, Dict

import numpy as np
import oracledb
import oci

# ------------- CONFIG -------------
EMBED_DIM = 1536
SKIP_OCI = os.getenv("SKIP_OCI") == "1"  # set SKIP_OCI=1 to bypass Object Storage

# ------------- OCI OBJECT STORAGE -------------
def make_os_client():
    """
    Builds an Object Storage client using ~/.oci/config (DEFAULT) if present,
    otherwise falls back to resource principals (e.g., Cloud Shell).
    """
    cfg_path = Path("~/.oci/config").expanduser()
    if cfg_path.exists():
        config = oci.config.from_file(str(cfg_path), "DEFAULT")
        return oci.object_storage.ObjectStorageClient(config)
    signer = oci.auth.signers.get_resource_principals_signer()
    return oci.object_storage.ObjectStorageClient(config={}, signer=signer)

def download_bucket_to(dir_path: Path, bucket: str) -> None:
    if SKIP_OCI:
        print("[skip] OCI download disabled via SKIP_OCI=1")
        return
    os_client = make_os_client()
    namespace = os_client.get_namespace().data
    objs = os_client.list_objects(namespace, bucket).data.objects
    for obj in objs:
        dest = dir_path / obj.name
        dest.parent.mkdir(parents=True, exist_ok=True)
        print(f"↓ {obj.name}")
        with open(dest, "wb") as f:
            data = os_client.get_object(namespace, bucket, obj.name).data
            f.write(data.content)

# ------------- SIMPLE PARSERS -------------
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import chardet

def read_txt(path: Path) -> Tuple[str, dict]:
    raw = path.read_bytes()
    enc = chardet.detect(raw).get("encoding") or "utf-8"
    return raw.decode(enc, errors="ignore"), {}

def read_pdf(path: Path) -> Tuple[str, dict]:
    reader = PdfReader(str(path))
    pages = []
    for i, pg in enumerate(reader.pages):
        try:
            t = pg.extract_text() or ""
        except Exception:
            t = ""
        pages.append(t)
    return "\n\n".join(pages), {"page_count": len(pages)}

def read_docx(path: Path) -> Tuple[str, dict]:
    doc = DocxDocument(str(path))
    paras = [p.text for p in doc.paragraphs if p.text]
    return "\n\n".join(paras), {}

def read_pptx(path: Path) -> Tuple[str, dict]:
    prs = Presentation(str(path))
    slides_txt: List[str] = []
    for s in prs.slides:
        parts = []
        for shp in s.shapes:
            if hasattr(shp, "text"):
                parts.append(shp.text)
        slides_txt.append("\n".join(parts))
    return "\n\n".join(slides_txt), {"slide_count": len(prs.slides)}

def load_text(path: Path) -> Tuple[str, dict]:
    ext = path.suffix.lower()
    if ext == ".pdf":  return read_pdf(path)
    if ext == ".docx": return read_docx(path)
    if ext == ".pptx": return read_pptx(path)
    if ext in {".txt", ".md"}: return read_txt(path)
    raise ValueError(f"Unsupported file type: {ext}")

# ------------- CHUNKING -------------
def window_chunks(text: str, max_chars=6000, overlap=700) -> List[str]:
    if not text: return []
    out, start, n = [], 0, len(text)
    while start < n:
        end = min(start + max_chars, n)
        out.append(text[start:end])
        if end == n: break
        start = max(0, end - overlap)
    return out

def header_aware_chunks(text: str, max_chars=6000, overlap=700) -> List[str]:
    import re
    if not text: return []
    paras = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]

    def is_header(p: str) -> bool:
        line = p.splitlines()[0] if p else ""
        if len(line) <= 80 and (line.isupper() or line.endswith(":")): return True
        return bool(re.match(r"^(\d+(\.\d+)*|[IVXLC]+\.)\s+\S", line))

    chunks, buf = [], ""
    for p in paras:
        add = ("\n\n" + p) if buf else p
        if len(buf + add) <= max_chars:
            buf = buf + add
        else:
            if buf: chunks.append(buf)
            tail = buf[-overlap:] if overlap and buf else ""
            buf = (tail + "\n\n" + p) if tail else p
        if is_header(p) and len(buf) > max_chars:
            chunks.append(buf); buf = ""
    if buf: chunks.append(buf)

    final = []
    for c in chunks:
        final.extend(window_chunks(c, max_chars, overlap) if len(c) > max_chars else [c])
    return final

def make_chunks(text: str, strategy="hybrid", max_chars=6000, overlap=700) -> List[str]:
    return header_aware_chunks(text, max_chars, overlap) if strategy != "window" else window_chunks(text, max_chars, overlap)

# ------------- EMBEDDINGS (PLACEHOLDER) -------------
def embed_texts(texts: List[str]) -> List[List[float]]:
    """
    Deterministic placeholder embeddings so pipeline works end-to-end.
    Swap with OCI GenAI embeddings later.
    """
    vecs = []
    for t in texts:
        seed = int(hashlib.md5(t.encode("utf-8", errors="ignore")).hexdigest(), 16) % (2**32)
        rng = np.random.RandomState(seed)
        v = rng.rand(EMBED_DIM).astype(float)
        v = v / np.linalg.norm(v)
        vecs.append(v.tolist())
    return vecs

# ------------- ADB CONNECTION -------------
def connect_adb():
    user = os.getenv("ADB_USER")
    pwd  = os.getenv("ADB_PASSWORD")
    dsn  = os.getenv("ADB_DSN")
    cfg  = os.getenv("TNS_ADMIN")
    missing = [k for k,v in [("ADB_USER",user),("ADB_PASSWORD",pwd),
                             ("ADB_DSN",dsn),("TNS_ADMIN",cfg)] if not v]
    if missing:
        raise RuntimeError(
            "Missing env vars: " + ", ".join(missing) + "\n"
            'Set them then re-run. Example:\n'
            'export TNS_ADMIN="/Users/rachelkim/Downloads/Wallet_chunking"\n'
            'export ADB_DSN="chunking_high"\n'
            'export ADB_USER="ADMIN"\n'
            'export ADB_PASSWORD="<your-admin-password>"'
        )
    return oracledb.connect(
        user=user,
        password=pwd,
        dsn=dsn,
        config_dir=cfg,
        wallet_location=cfg,
        ssl_server_dn_match=True
    )

# ------------- SCHEMA (auto-create if missing) -------------
def init_schema(conn):
    ddl_table = """
    CREATE TABLE doc_chunks (
      id         NUMBER GENERATED BY DEFAULT AS IDENTITY PRIMARY KEY,
      course     VARCHAR2(64),
      source_uri VARCHAR2(1024),
      chunk_no   NUMBER,
      section    VARCHAR2(256),
      page_from  NUMBER,
      page_to    NUMBER,
      content    CLOB,
      embedding  VECTOR(1536),
      metadata   JSON
    )
    """
    ddl_idx1 = "CREATE INDEX doc_chunks_course_idx ON doc_chunks(course)"
    ddl_idx2 = "CREATE INDEX doc_chunks_source_idx ON doc_chunks(source_uri)"

    ddl_vec_new = """
    CREATE VECTOR INDEX doc_chunks_vec_idx
    ON doc_chunks (embedding)
    PARAMETERS('ALGORITHM HNSW DISTANCE COSINE')
    """
    ddl_vec_legacy = """
    CREATE INDEX doc_chunks_vec_idx
    ON doc_chunks (embedding)
    ORGANIZATION NEIGHBOR LIST
    """
    ddl_vec_alter = """
    ALTER INDEX doc_chunks_vec_idx
    PARAMETERS('ALGORITHM HNSW DISTANCE COSINE')
    """

    with conn.cursor() as cur:
        # table
        cur.execute("SELECT COUNT(*) FROM user_tables WHERE table_name='DOC_CHUNKS'")
        if cur.fetchone()[0] == 0:
            print("[schema] creating table DOC_CHUNKS")
            cur.execute(ddl_table)
            conn.commit()
        else:
            print("[schema] DOC_CHUNKS already exists")

        # conventional indexes (ignore if exist)
        for ddl in (ddl_idx1, ddl_idx2):
            try:
                cur.execute(ddl); conn.commit()
            except oracledb.Error:
                pass

        # vector index (new syntax, then legacy)
        cur.execute("SELECT COUNT(*) FROM user_indexes WHERE index_name='DOC_CHUNKS_VEC_IDX'")
        if cur.fetchone()[0] == 0:
            created = False
            try:
                print("[schema] creating vector index (new syntax)")
                cur.execute(ddl_vec_new); conn.commit(); created = True
            except oracledb.Error:
                try:
                    print("[schema] creating vector index (legacy syntax)")
                    cur.execute(ddl_vec_legacy); conn.commit(); created = True
                except oracledb.Error:
                    print("[schema] could not create vector index; continue without (queries still work, just slower)")
            if created:
                try:
                    cur.execute(ddl_vec_alter); conn.commit()
                except oracledb.Error:
                    pass
        else:
            print("[schema] vector index already exists")

# ------------- INSERTS -------------
def insert_chunks(conn, rows: List[Dict]):
    import json as _json
    # normalize types
    for r in rows:
        r["embedding"] = [float(x) for x in r["embedding"]]
        r["metadata"]  = _json.dumps(r.get("metadata", {}))

    sql = """
    insert into doc_chunks
      (course, source_uri, chunk_no, section, page_from, page_to, content, embedding, metadata)
    values
      (:course, :source_uri, :chunk_no, :section, :page_from, :page_to, :content, :embedding, :metadata)
    """
    with conn.cursor() as cur:
        cur.setinputsizes(
            course=oracledb.DB_TYPE_VARCHAR,
            source_uri=oracledb.DB_TYPE_VARCHAR,
            chunk_no=oracledb.DB_TYPE_NUMBER,
            section=oracledb.DB_TYPE_VARCHAR,
            page_from=oracledb.DB_TYPE_NUMBER,
            page_to=oracledb.DB_TYPE_NUMBER,
            content=oracledb.DB_TYPE_CLOB,
            embedding=oracledb.DB_TYPE_VECTOR,
            metadata=oracledb.DB_TYPE_JSON
        )
        cur.executemany(sql, rows)
        conn.commit()

# ------------- MAIN -------------
def main():
    COURSE = os.environ.get("COURSE", "DATA100_Fa25")
    BUCKET = os.environ.get("BUCKET", "course-docs")
    RAWDIR = Path("rag_course/docs/raw")
    RAWDIR.mkdir(parents=True, exist_ok=True)

    # 1) pull files (or skip)
    download_bucket_to(RAWDIR, BUCKET)

    # 2) parse + chunk
    rows_all: List[Dict] = []
    for p in sorted(RAWDIR.rglob("*")):
        if not p.is_file(): continue
        if p.suffix.lower() not in {".pdf",".docx",".pptx",".txt",".md"}: continue
        print(f"[parse] {p.name}")
        text, meta = load_text(p)
        chunks = make_chunks(text, strategy="hybrid", max_chars=6000, overlap=700)
        page_to = meta.get("page_count") or meta.get("slide_count")
        for i, c in enumerate(chunks, start=1):
            rows_all.append({
                "course": COURSE,
                "source_uri": f"oci://{BUCKET}/{p.name}" if not SKIP_OCI else f"local://{p.name}",
                "chunk_no": i,
                "section": None,
                "page_from": 1 if page_to else None,
                "page_to": page_to,
                "content": c,
                "metadata": {"strategy":"hybrid","bucket":BUCKET,"file":p.name}
            })
    print(f"[chunks] total: {len(rows_all)}")

    # 3) embed (placeholder)
    BATCH = 32
    for i in range(0, len(rows_all), BATCH):
        batch = rows_all[i:i+BATCH]
        vecs = embed_texts([r["content"] for r in batch])
        for r, v in zip(batch, vecs):
            r["embedding"] = v

    # 4) insert into ADB
    conn = connect_adb()
    init_schema(conn)
    insert_chunks(conn, rows_all)
    print("[done] inserted rows into doc_chunks")

if __name__ == "__main__":
    main()
    